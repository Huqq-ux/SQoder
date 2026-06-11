import os
import tempfile
from pathlib import Path
from Coder.knowledge.document_loader import DocumentLoader

_PROJECT_DIR = os.path.normpath(
    os.path.join(os.path.dirname(__file__), "..", "..")
)


class TestDocumentLoaderNewFormats:
    def setup_method(self):
        self.loader = DocumentLoader()
        self._tmpdir = tempfile.mkdtemp(dir=_PROJECT_DIR)

    def teardown_method(self):
        import shutil
        if os.path.exists(self._tmpdir):
            shutil.rmtree(self._tmpdir)

    def _make_temp(self, suffix: str) -> str:
        fd, path = tempfile.mkstemp(suffix=suffix, dir=self._tmpdir)
        os.close(fd)
        return path

    def test_load_pptx(self):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "第一章 绪论"
        body = slide.shapes.placeholders[1]
        body.text = "这是课程介绍内容"

        path = self._make_temp(".pptx")
        try:
            prs.save(path)
            result = self.loader.load(path)
            assert "第一章 绪论" in result["content"]
            assert "课程介绍内容" in result["content"]
        finally:
            os.unlink(path)

    def test_load_xlsx(self):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.append(["题目", "选项A", "选项B", "答案"])
        ws.append(["1+1=?", "1", "2", "B"])

        path = self._make_temp(".xlsx")
        try:
            wb.save(path)
            result = self.loader.load(path)
            assert "1+1=?" in result["content"]
            assert "题目" in result["content"]
        finally:
            os.unlink(path)

    def test_load_csv(self):
        import csv as csv_module
        path = self._make_temp(".csv")
        try:
            with open(path, "w", encoding="utf-8", newline="") as f:
                writer = csv_module.writer(f)
                writer.writerow(["题目", "答案"])
                writer.writerow(["极限的定义", "见教材P12"])
            result = self.loader.load(path)
            assert "极限的定义" in result["content"]
        finally:
            os.unlink(path)
